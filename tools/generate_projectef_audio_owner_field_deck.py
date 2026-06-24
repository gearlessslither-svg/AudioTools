from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORT_DIR = ROOT / "报告"
OUT = REPORT_DIR / "ProjectEF_音频Owner与开工字段机制_2026-06-16.pptx"

FONT = "Microsoft YaHei"
FONT_BOLD = "Microsoft YaHei UI"

COLORS = {
    "bg": RGBColor(247, 248, 250),
    "panel": RGBColor(255, 255, 255),
    "ink": RGBColor(28, 36, 46),
    "muted": RGBColor(96, 108, 124),
    "line": RGBColor(218, 224, 232),
    "teal": RGBColor(18, 122, 116),
    "blue": RGBColor(56, 101, 177),
    "amber": RGBColor(222, 142, 50),
    "red": RGBColor(198, 75, 75),
    "green": RGBColor(72, 148, 90),
    "purple": RGBColor(125, 88, 166),
    "slate": RGBColor(63, 78, 96),
}


def hex_color(s: str) -> RGBColor:
    s = s.strip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]


def style_run(run, size=16, color=None, bold=False):
    run.font.name = FONT_BOLD if bold else FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def textbox(slide, x, y, w, h, text, size=16, color=None, bold=False, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    style_run(r, size=size, color=color, bold=bold)
    return shape


def title(slide, text, kicker=None):
    if kicker:
        textbox(slide, 0.62, 0.36, 7.8, 0.28, kicker, size=10.5, color=COLORS["teal"], bold=True)
    textbox(slide, 0.60, 0.66, 11.6, 0.66, text, size=27, color=COLORS["ink"], bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.64), Inches(1.36), Inches(1.16), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["teal"]
    line.line.fill.background()


def footer(slide, page):
    textbox(slide, 0.62, 7.05, 6.2, 0.22, "ProjectEF Audio Owner & Field Contract", size=8.5, color=COLORS["muted"])
    textbox(slide, 11.94, 7.05, 0.70, 0.22, f"{page:02d}", size=8.5, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, heading, body, color, heading_size=14, body_size=10.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["panel"]
    shape.line.color.rgb = COLORS["line"]
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    textbox(slide, x + 0.18, y + 0.14, w - 0.32, 0.34, heading, size=heading_size, color=color, bold=True)
    textbox(slide, x + 0.18, y + 0.54, w - 0.34, h - 0.64, body, size=body_size, color=COLORS["slate"])


def bullets(slide, x, y, w, h, items, size=14):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = Pt(5)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["ink"]


def table(slide, x, y, w, h, headers, rows, widths=None, font_size=9.5, header_color=None):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    if widths:
        for i, width in enumerate(widths):
            tbl.columns[i].width = Inches(width)
    for j, head in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color or COLORS["teal"]
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        p = cell.text_frame.paragraphs[0]
        p.text = head
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = FONT_BOLD
        p.runs[0].font.size = Pt(font_size)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = RGBColor(255, 255, 255)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if i % 2 else RGBColor(250, 252, 254)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            p.runs[0].font.name = FONT
            p.runs[0].font.size = Pt(font_size)
            p.runs[0].font.color.rgb = COLORS["ink"]


def flow(slide, items, y=2.0):
    x = 0.82
    w = 1.58
    gap = 0.22
    for i, (head, body, color) in enumerate(items):
        card(slide, x + i * (w + gap), y, w, 1.36, head, body, color, heading_size=11.5, body_size=8.4)
        if i < len(items) - 1:
            textbox(slide, x + i * (w + gap) + w + 0.02, y + 0.50, 0.22, 0.22, "→", size=16, color=COLORS["muted"], bold=True)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    page = 1

    # 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_color("F3F7F6")
    bg.line.fill.background()
    textbox(slide, 0.86, 0.82, 9.6, 0.30, "ProjectEF Audio Production Contract", size=12, color=COLORS["teal"], bold=True)
    textbox(slide, 0.84, 1.40, 11.6, 1.04, "音频 Owner 与开工字段机制", size=38, color=COLORS["ink"], bold=True)
    textbox(slide, 0.88, 2.62, 10.8, 0.60, "让策划、美术、程序、QA 在正确节点给音频提供必要开工信息", size=21, color=COLORS["slate"], bold=True)
    textbox(slide, 0.90, 4.02, 10.4, 0.72, "核心定义：音频开工信息 = 策划可感知需求 + 美术资源映射 + 程序触发合约 + QA 测试路径", size=19, color=COLORS["teal"], bold=True)
    textbox(slide, 0.92, 6.62, 3.8, 0.24, "生成日期：2026-06-16", size=9.5, color=COLORS["muted"])
    footer(slide, page)
    page += 1

    # 2 thesis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "核心结论", "Executive Summary")
    textbox(slide, 0.92, 1.60, 11.45, 0.56, "音频不是缺执行力，而是缺被上游结构化提供的开工信息。", size=22, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    card(slide, 0.84, 2.48, 3.72, 2.00, "策划负责", "提出玩家可感知、可判断、可验收的声音需求。策划不写 Wwise，但要说明玩家应该听懂什么。", COLORS["blue"])
    card(slide, 4.82, 2.48, 3.72, 2.00, "资源/逻辑负责", "美术/动画/VFX/程序提供资源名、路径、状态、触发和参数。音频不应靠猜找资源。", COLORS["amber"])
    card(slide, 8.80, 2.48, 3.72, 2.00, "测试负责", "QA/程序/策划提供测试入口、目标和验收方式。没有测试路径，音频无法闭环。", COLORS["red"])
    textbox(slide, 1.02, 5.24, 11.25, 0.62, "制度化规则：每个音频需求都必须有需求 owner、资源/逻辑 owner、测试 owner；缺任意一个，只能进入音频探索，不能进入正式交付。", size=18, color=COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 3 problem example shortcut
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "例 1：快捷键功能为什么缺失音效", "Design Document Responsibility")
    card(slide, 0.86, 1.58, 3.55, 1.82, "策划问题", "快捷键功能为什么没有音效？玩家按下快捷键时应该有反馈。", COLORS["blue"])
    card(slide, 4.90, 1.58, 3.55, 1.82, "音频反问", "什么是快捷键功能？需求在哪里？成功/失败状态是什么？怎么测试？", COLORS["red"])
    card(slide, 8.94, 1.58, 3.55, 1.82, "流程结论", "如果策划可以判断这里缺声音，策划文档或 Jira 就应该记录该可感知需求。", COLORS["teal"])
    rows = [
        ["功能名", "快捷键切换/使用"],
        ["玩家行为", "按下快捷键触发功能"],
        ["预期反馈", "成功操作有确认音；不可用/失败状态有区别反馈"],
        ["验收方式", "策划测试快捷键时能判断有无反馈、反馈是否符合成功/失败"],
    ]
    table(slide, 1.22, 4.22, 10.88, 1.72, ["策划文档至少记录", "示例"], rows, widths=[2.1, 8.68], font_size=10.5)
    footer(slide, page)
    page += 1

    # 4 design responsibility
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "策划文档中的音频需求：分两层写", "Design Requirement Layers")
    card(slide, 0.92, 1.62, 5.45, 3.10, "策划初始需求", "说明这个功能/状态是否需要玩家听到反馈；成功、失败、变化、危险、奖励、提示分别是什么；作为策划可验收的标准。", COLORS["blue"], heading_size=16, body_size=13)
    card(slide, 6.92, 1.62, 5.45, 3.10, "音频补全需求", "音频判断用 UI、SFX、环境、状态 Loop、RTPC、Switch、音乐 Sting 还是不做；补资源量、Wwise 合约、混音和测试证据。", COLORS["teal"], heading_size=16, body_size=13)
    textbox(slide, 1.04, 5.45, 11.25, 0.50, "真正完备的音频需求 = 策划提出可感知需求 + 音频拆解实现需求 + 测试后补齐验收需求", size=18, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.22, 6.14, 10.88, 0.36, "如果通过优化单后补，也可以，但要明确这是版本后发现需求，会带来全貌不可预判和内容缺失风险。", size=13.5, color=COLORS["red"], align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 5 art example fish splash
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "例 2：鱼炸水需要音效，但怎么配？", "Art Resource Map")
    card(slide, 0.84, 1.55, 3.75, 1.68, "策划描述", "鱼炸水需要音效。", COLORS["blue"])
    card(slide, 4.80, 1.55, 3.75, 1.68, "音频问题", "有 8 个 prefab，水声有大中小三档，应该由资源名区分还是程序参数区分？", COLORS["red"])
    card(slide, 8.76, 1.55, 3.75, 1.68, "当前断点", "策划只知道一个 prefab，无法覆盖全部资源和表现关系。", COLORS["amber"])
    rows = [
        ["策划名", "鱼炸水 / 鱼跃水 / 入水水花"],
        ["Unity 资源名", "Prefab / Animation / Timeline / VFX 名"],
        ["资源路径", "可定位到项目文件"],
        ["表现等级", "大/中/小、强/弱、普通/稀有"],
        ["触发关系", "谁触发它，是否和程序状态绑定"],
        ["音频建议字段", "AudioEventId / AudioProfileId / RTPC 参数"],
        ["Owner", "美术 / TA / 动画 / VFX / 程序负责人"],
    ]
    table(slide, 1.02, 3.66, 11.32, 2.38, ["Art Resource Map 字段", "说明"], rows, widths=[2.1, 9.1], font_size=9.8)
    footer(slide, page)
    page += 1

    # 6 upstream resource list
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "美术/资源流程应前置提供哪些文件名", "Resource Fields")
    rows = [
        ["Prefab", "路径、挂点、组件、是否运行时生成、owner、是否可加音频字段"],
        ["Animation", "clip 名、关键帧、循环段、是否可加 Animation Event、预览方式"],
        ["Timeline", "Timeline 名、Track、Marker、播放时机、是否复用"],
        ["VFX", "VFX 名、播放时机、强弱等级、持续时间、是否跟随对象"],
        ["UI Prefab", "按钮组件、pressed/hover/success/fail 字段、默认音频规则"],
        ["Config", "表名、ID、AudioEventId/AudioProfileId 字段、是否热更"],
        ["State Machine", "状态名、进入/退出条件、状态切换频率、owner"],
    ]
    table(slide, 0.94, 1.56, 11.46, 4.68, ["资源类型", "音频需要的信息"], rows, widths=[1.75, 9.58], font_size=10.3)
    textbox(slide, 1.08, 6.42, 11.10, 0.28, "名单可以逐渐加长；当前至少需要 Prefab、Animation、Timeline、VFX、UI Prefab、Config。", size=13.2, color=COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 7 Jira phases
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Jira 流转建议：把音频拆成三个阶段", "Jira Flow")
    flow(slide, [
        ("Audio Explore\nReady", "资源确认\n功能未跑通", COLORS["blue"]),
        ("Audio Production\nReady", "策划可测\n音频可做", COLORS["green"]),
        ("Audio Acceptance\nReady", "系统正常\n音频验收", COLORS["teal"]),
        ("Audio Done", "策划验收\n证据归档", COLORS["slate"]),
    ], y=1.76)
    rows = [
        ["Explore Ready", "美术资源基本确认，但程序逻辑跑不起来；音频可以前置判断资源、素材方向、Wwise 模板、命名和参数方案。"],
        ["Production Ready", "系统功能可以交给策划测试；音频第一时间收到测试方式、测试目标、流程清单、Debug 入口，开始正式制作和配置。"],
        ["Acceptance Ready", "策划测试完所有正常表现并确认系统无问题；音频进行专项测试、修正配置，再交策划验收声音。"],
    ]
    table(slide, 1.00, 4.02, 11.30, 1.78, ["阶段", "含义"], rows, widths=[2.0, 9.15], font_size=10.4)
    textbox(slide, 1.06, 6.16, 11.18, 0.34, "关键：可前置探索 ≠ 可最终交付。不同阶段承诺不同，不混用 Ready。", size=15, color=COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 8 stage entry conditions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "三个阶段的进入条件", "Entry Criteria")
    rows = [
        ["Audio Explore Ready", "系统大概玩法已知；美术资源名/路径初步存在；资源表现等级大概清楚；程序逻辑可未完成或不可测。"],
        ["Audio Production Ready", "策划能跑通功能；测试路径明确；触发方式明确；程序/美术资源基本稳定；音频配置点明确。"],
        ["Audio Acceptance Ready", "系统逻辑已验收；音频已接入；QA/音频可复现；策划能验收听感和反馈是否正确。"],
    ]
    table(slide, 0.92, 1.58, 11.55, 3.10, ["阶段", "进入条件"], rows, widths=[2.45, 8.95], font_size=11.2)
    card(slide, 1.10, 5.22, 10.98, 0.88, "状态变化通知", "策划能测系统的时候，音频必须第一时间知道；否则音频无法完整感知系统内容，也无法及时制作和配置。", COLORS["amber"], heading_size=13, body_size=11.5)
    footer(slide, page)
    page += 1

    # 9 notify responsibility
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "新增字段：Audio Notify Owner", "Notification Contract")
    rows = [
        ["Audio Notify Owner", "谁负责通知音频这个系统可测"],
        ["Test Available Date", "什么时候可以开始测试"],
        ["Test Entry", "场景、入口、账号、存档、Debug 命令"],
        ["Test Scope", "这次能测哪些正常表现和异常表现"],
        ["Known Limitations", "哪些还不能测、哪些是临时表现"],
        ["Evidence Link", "测试文档、视频、截图、日志、Profiler 报告位置"],
    ]
    table(slide, 1.00, 1.58, 11.20, 3.85, ["字段", "说明"], rows, widths=[2.55, 8.55], font_size=11.0)
    textbox(slide, 1.12, 5.92, 11.05, 0.44, "这不是细节，这是流程生命线：系统可测但音频不知道，等同于音频被动漏需求。", size=17, color=COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 10 owner triad
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "每个音频需求至少有三个 Owner", "Owner Triad")
    card(slide, 0.86, 1.64, 3.58, 2.34, "需求 Owner", "通常是策划。负责说明玩家需要听懂什么、什么状态需要反馈、成功/失败/奖励/危险如何验收。", COLORS["blue"], heading_size=16, body_size=12)
    card(slide, 4.88, 1.64, 3.58, 2.34, "资源/逻辑 Owner", "美术、动画、VFX、程序。负责资源名、路径、状态机、触发函数、参数范围和生命周期。", COLORS["amber"], heading_size=16, body_size=12)
    card(slide, 8.90, 1.64, 3.58, 2.34, "测试 Owner", "QA、程序或策划。负责测试入口、目标、步骤、账号/存档、Debug 和验收证据。", COLORS["red"], heading_size=16, body_size=12)
    textbox(slide, 1.12, 4.72, 11.05, 0.50, "缺任意一个 Owner：只能进入 Audio Explore，不能进入 Audio Production Ready。", size=19, color=COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 11 acceptance boundary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "策划验收、音频验收、QA 验收的边界", "Acceptance Boundary")
    rows = [
        ["策划验收", "状态有没有声音反馈；成功/失败/危险/奖励是否听得懂；声音是否符合功能意图；是否干扰核心玩法判断。"],
        ["音频验收", "Event 是否正确触发/停止；RTPC/Switch/State 是否正确；混音、优先级、随机、衰减、多人裁剪是否正确。"],
        ["QA 验收", "测试步骤是否覆盖；正常/异常/边界是否可复现；回归是否稳定；证据是否可追溯。"],
    ]
    table(slide, 0.98, 1.62, 11.36, 2.90, ["验收角色", "验收内容"], rows, widths=[1.72, 9.50], font_size=11.0)
    card(slide, 1.04, 5.08, 11.26, 0.88, "边界原则", "策划验收音频体验，不等于判断 Wwise 是否正确；音频验收实现正确性，不等于代替 QA 跑全场景回归。", COLORS["teal"], heading_size=13, body_size=11.2)
    footer(slide, page)
    page += 1

    # 12 Jira fields
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "建议加入 Jira / 需求表的字段", "Field Template")
    rows = [
        ["Audio Required?", "Yes / No / Unknown；Unknown 自动进入音频 Triage"],
        ["Audio Perception Need", "玩家需要听懂什么信息：操作、状态、危险、奖励、空间、情绪？"],
        ["Design Acceptance", "策划如何判断声音反馈是对的"],
        ["Art Resource Map", "Prefab / Animation / Timeline / VFX / UI / Config 路径与 owner"],
        ["Program Trigger Contract", "触发 owner、Play/Stop、RTPC、Switch/State、生命周期、Debug"],
        ["Audio Notify Owner", "谁在系统可测时通知音频"],
        ["QA Recipe", "场景、步骤、账号/存档、Debug、预期、证据"],
        ["Audio Stage", "Candidate / Explore Ready / Production Ready / Acceptance Ready / Done"],
    ]
    table(slide, 0.82, 1.50, 11.85, 5.05, ["字段", "说明"], rows, widths=[2.62, 9.08], font_size=9.8)
    footer(slide, page)
    page += 1

    # 13 change regression
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "例 3：上个版本没问题，这个版本为什么又出问题？", "Change Regression")
    rows = [
        ["脚步声材质 / 快速走路", "上个版本已经测试完毕并调整好；这个版本又出现问题。"],
        ["主界面 UI 音效", "上个版本已经配置好；这个版本没有了。"],
        ["营地的火", "上个版本已经配置好；这个版本消失了。"],
        ["环境地图可行走区域", "可行走区域改变，但没有人告知音频；需要音频想到后自行询问。"],
    ]
    table(slide, 0.86, 1.52, 11.65, 2.55, ["现象", "版本回归表现"], rows, widths=[3.15, 8.35], font_size=11.0, header_color=COLORS["red"])
    card(
        slide,
        0.92,
        4.34,
        5.55,
        1.46,
        "底层问题",
        "音频不知道上游进行了哪些更改；上游也不知道哪些更改会影响音频。因此旧版本已完成的声音，在新版本里仍可能静默失效。",
        COLORS["red"],
        heading_size=13,
        body_size=11.2,
    )
    card(
        slide,
        6.86,
        4.34,
        5.55,
        1.46,
        "解决思路",
        "任何有音效需求单的 Jira，只要上游进行了设计、美术、程序、场景或配置修改，都需要派生/重开新的“音效确认”单，进入音频回归确认。",
        COLORS["teal"],
        heading_size=13,
        body_size=11.2,
    )
    textbox(slide, 1.10, 6.18, 11.15, 0.46, "规则不是要求上游判断 Wwise 细节，而是让变更 Owner 告诉音频：这个系统变了，需要你确认是否影响声音。", size=15.5, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 14 progress visibility
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "音频进度需要纳入研发整体进度", "Progress Visibility")
    card(
        slide,
        0.86,
        1.46,
        11.62,
        0.92,
        "当前问题",
        "项目会议和项目文档体现不出音频进度，导致音频需求不可视、不可控；团队会习惯性忽视音频工作量的堆积，最后只能靠后期人力强推。",
        COLORS["red"],
        heading_size=13,
        body_size=11.2,
    )
    rows = [
        ["音频功能开发进度", "看音频相关功能是否具备开工/接入条件", "Wwise/Unity 接口、触发逻辑、配置字段、Debug 工具、测试入口、阻塞 Owner"],
        ["音频需求完成进度", "看每个系统的声音需求是否完成闭环", "需求数、已制作、已配置、已测试、已验收、待确认、返工/回归风险"],
    ]
    table(slide, 0.86, 2.74, 11.62, 1.74, ["维度", "管理目标", "会议/文档需要展示"], rows, widths=[2.15, 3.20, 6.12], font_size=10.5, header_color=COLORS["teal"])
    card(
        slide,
        0.92,
        4.82,
        5.55,
        1.22,
        "会议机制",
        "项目例会中固定出现 Audio Status：本周完成、当前阻塞、下周风险、需要上游决策/资源/测试支持。",
        COLORS["blue"],
        heading_size=13,
        body_size=11.0,
    )
    card(
        slide,
        6.86,
        4.82,
        5.55,
        1.22,
        "文档机制",
        "项目进度表 / Jira Dashboard 增加音频列：Audio Stage、Audio Owner、Dependency Owner、Test Ready、Audio Acceptance。",
        COLORS["purple"],
        heading_size=13,
        body_size=11.0,
    )
    textbox(slide, 1.00, 6.38, 11.34, 0.40, "原则：功能没有音频状态，就不能被视为研发进度完全透明。", size=17, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 15 final flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "推荐最终流程", "End-to-End")
    flow(slide, [
        ("策划文档", "可感知\n音频需求", COLORS["blue"]),
        ("资源映射", "Prefab/Anim\nVFX/Config", COLORS["amber"]),
        ("程序合约", "触发/参数\nDebug", COLORS["purple"]),
        ("可测通知", "入口/目标\n限制", COLORS["green"]),
        ("音频制作", "资源/Wwise\nUnity 配置", COLORS["teal"]),
        ("音频测试", "日志/Profiler\n修正", COLORS["red"]),
        ("策划验收", "体验正确\n任务关闭", COLORS["slate"]),
    ], y=1.92)
    textbox(slide, 1.00, 4.20, 11.32, 0.58, "音频的工作不再从“你们觉得缺什么声音”开始，而是从字段完整的任务开始。", size=19, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.05, 5.12, 11.18, 0.66, "这样音频可以前置探索、稳定制作、可复现测试、可追踪验收；上游也知道自己需要给音频什么。", size=16, color=COLORS["slate"], align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 16 close
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    textbox(slide, 0.92, 1.18, 11.48, 0.78, "一句话落地", size=34, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.10, 2.34, 11.10, 1.20, "从每一个环节的 Owner 开始，找到负责人，补齐字段，再让音频开工。", size=30, color=COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.32, 4.12, 10.70, 0.78, "缺需求 owner、资源/逻辑 owner、测试 owner 的任务，只能探索，不能正式交付。", size=20, color=COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)

    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    candidates = [OUT] + [OUT.with_name(f"{OUT.stem}_v{i}{OUT.suffix}") for i in range(2, 10)]
    last_error = None
    for final_out in candidates:
        try:
            prs.save(final_out)
        except PermissionError as exc:
            last_error = exc
            continue
        print(f"PPTX: {final_out}")
        break
    else:
        raise last_error


if __name__ == "__main__":
    build()
