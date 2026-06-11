# -*- coding: utf-8 -*-
from __future__ import annotations

import copy
import mimetypes
import os
import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\周报\2026 0504 - 0515 -  双周报")
SOURCE_PPT = ROOT / "音频中心工作汇报 0504 - 0515.pptx"
OUT_PPT = Path(r"G:\AI\Material\Wwise\音频中心工作汇报 0504 - 0515 - 已整理.pptx")
POSTER_DIR = Path(r"G:\AI\Material\Wwise\_weekly_ppt_posters")


def read_docx(path: Path) -> str:
    ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    with zipfile.ZipFile(path) as z:
        xml = z.read("word/document.xml")
    root = ET.fromstring(xml)
    out: list[str] = []
    for p in root.iter(ns + "p"):
        text = "".join((t.text or "") for t in p.iter(ns + "t")).strip()
        if text:
            out.append(text)
    return "\n".join(out)


def read_binary_doc_text(path: Path) -> list[str]:
    # The old .doc file stores the useful body text as UTF-16LE fragments.
    b = path.read_bytes()
    hits: list[str] = []
    cur: list[str] = []
    allowed = "，。；：、（）《》【】！？\n\r\t /*-"
    for i in range(0, len(b) - 1, 2):
        code = b[i] + (b[i + 1] << 8)
        ch = chr(code)
        if (0x20 <= code <= 0x7E) or (0x4E00 <= code <= 0x9FFF) or ch in allowed:
            cur.append(ch)
            continue
        if len(cur) >= 4:
            s = "".join(cur).strip()
            if s and any("\u4e00" <= c <= "\u9fff" for c in s):
                hits.append(clean_doc_fragment(s))
        cur = []
    if len(cur) >= 4:
        s = "".join(cur).strip()
        if s and any("\u4e00" <= c <= "\u9fff" for c in s):
            hits.append(clean_doc_fragment(s))
    noise = {"鮥螑獽彩"}
    return [h for h in hits if h and h not in noise and "愀氀" not in h]


def clean_doc_fragment(text: str) -> str:
    text = text.replace("\t", " ")
    text = re.sub(r"[ \u3000]+", " ", text)
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def set_text(shape, text: str, font_size: int | None = None, bold_first: bool = False) -> None:
    shape.text = text
    tf = shape.text_frame
    tf.word_wrap = True
    for pi, p in enumerate(tf.paragraphs):
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            if font_size is not None:
                run.font.size = Pt(font_size)
            run.font.name = "Microsoft YaHei"
            run.font.color.rgb = RGBColor(40, 40, 40)
            if bold_first and pi == 0:
                run.font.bold = True


def iter_text_shapes(slide):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            yield shape


def title_shape(slide):
    candidates = [s for s in iter_text_shapes(slide) if "标题" in s.name]
    if candidates:
        return candidates[0]
    candidates = sorted(iter_text_shapes(slide), key=lambda s: (s.top, s.left))
    return candidates[0] if candidates else None


def content_shape(slide):
    t_shape = title_shape(slide)
    candidates = [s for s in iter_text_shapes(slide) if "内容占位符" in s.name or "矩形" in s.name]
    if candidates:
        return sorted(candidates, key=lambda s: s.top)[-1]
    texts = [
        s
        for s in iter_text_shapes(slide)
        if t_shape is None or s.element is not t_shape.element
    ]
    return sorted(texts, key=lambda s: (s.top, s.left))[0] if texts else None


def set_title(slide, text: str) -> None:
    shape = title_shape(slide)
    if shape is None:
        return
    set_text(shape, text, font_size=30, bold_first=True)
    for p in shape.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.color.rgb = RGBColor(31, 78, 121)


def set_body(slide, text: str, font_size: int = 22) -> None:
    shape = content_shape(slide)
    if shape is None:
        return
    set_text(shape, text, font_size=font_size, bold_first=True)


def clear_body(slide) -> None:
    shape = content_shape(slide)
    if shape is not None:
        shape.text = ""


def duplicate_slide(prs: Presentation, index: int, insert_after: int | None = None):
    source = prs.slides[index]
    blank = prs.slide_layouts[6]
    dest = prs.slides.add_slide(blank)

    for shape in source.shapes:
        newel = copy.deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, "p:extLst")

    for rel in source.part.rels.values():
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype:
            continue
        dest.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)

    if insert_after is not None:
        sld_id_lst = prs.slides._sldIdLst
        new_id = sld_id_lst[-1]
        sld_id_lst.remove(new_id)
        sld_id_lst.insert(insert_after + 1, new_id)
    return dest


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    prs.slides._sldIdLst.remove(slide_id)


def font_path() -> str | None:
    for path in [
        r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\simhei.ttf",
        r"C:\Windows\Fonts\simsun.ttc",
        r"C:\Windows\Fonts\arial.ttf",
    ]:
        if Path(path).exists():
            return path
    return None


FONT_FILE = font_path()


def text_size(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont) -> tuple[int, int]:
    bbox = draw.textbbox((0, 0), text, font=font)
    return bbox[2] - bbox[0], bbox[3] - bbox[1]


def make_poster(label: str, out: Path, size: tuple[int, int] = (1280, 720)) -> Path:
    out.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", size, (235, 238, 242))
    draw = ImageDraw.Draw(img)
    w, h = size
    draw.rectangle((0, 0, w, h), fill=(235, 238, 242))
    draw.rectangle((0, h - 92, w, h), fill=(31, 78, 121))
    draw.polygon(
        [(w // 2 - 70, h // 2 - 85), (w // 2 - 70, h // 2 + 85), (w // 2 + 95, h // 2)],
        fill=(31, 78, 121),
    )
    font_big = ImageFont.truetype(FONT_FILE, 52) if FONT_FILE else ImageFont.load_default()
    font_small = ImageFont.truetype(FONT_FILE, 34) if FONT_FILE else ImageFont.load_default()
    label = label[:30]
    tw, th = text_size(draw, label, font_big)
    draw.text(((w - tw) / 2, h - 72), label, font=font_big, fill=(255, 255, 255))
    sub = "视频素材"
    sw, sh = text_size(draw, sub, font_small)
    draw.text(((w - sw) / 2, 58), sub, font=font_small, fill=(31, 78, 121))
    img.save(out)
    return out


def add_label(slide, text: str, left, top, width, height=Inches(0.28), font_size=12) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box, text, font_size=font_size)
    for p in box.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.color.rgb = RGBColor(31, 78, 121)


def add_video(slide, path: Path, left, top, width, height, label: str | None = None) -> None:
    label = label or path.stem
    poster = make_poster(label, POSTER_DIR / f"{safe_name(label)}.png")
    mime = mimetypes.guess_type(path.name)[0] or "video/unknown"
    slide.shapes.add_movie(str(path), left, top, width, height, poster_frame_image=str(poster), mime_type=mime)
    add_label(slide, label, left, top + height + Inches(0.03), width)


def safe_name(text: str) -> str:
    return re.sub(r"[^\w\u4e00-\u9fff-]+", "_", text)[:60]


def add_audio_list(slide, title: str, items: list[str], left, top, width, height) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    text = title + "\n" + "\n".join(f"{i + 1}. {item}" for i, item in enumerate(items))
    set_text(box, text, font_size=16, bold_first=True)
    for pi, p in enumerate(box.text_frame.paragraphs):
        for run in p.runs:
            run.font.color.rgb = RGBColor(31, 78, 121) if pi == 0 else RGBColor(50, 50, 50)


def add_grid_videos(slide, videos: list[Path], labels: list[str]) -> None:
    clear_body(slide)
    lefts = [Inches(0.68), Inches(6.82)]
    tops = [Inches(1.15), Inches(4.0)]
    w = Inches(5.65)
    h = Inches(2.35)
    for idx, path in enumerate(videos[:4]):
        row, col = divmod(idx, 2)
        add_video(slide, path, lefts[col], tops[row], w, h, labels[idx])


def main() -> None:
    prs = Presentation(str(SOURCE_PPT))

    # Add two extra LL output pages after slide 9, cloned from the existing LL output page.
    duplicate_slide(prs, 8, insert_after=8)
    duplicate_slide(prs, 8, insert_after=9)

    # Cover.
    slide = prs.slides[0]
    text_shapes = list(iter_text_shapes(slide))
    if len(text_shapes) >= 2:
        set_text(text_shapes[1], "音频中台 工作汇报", font_size=44, bold_first=True)
        for p in text_shapes[1].text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.color.rgb = RGBColor(31, 78, 121)
    if len(text_shapes) >= 3:
        set_text(text_shapes[2], "（5月4日-5月15日 双周报）", font_size=28)
        for p in text_shapes[2].text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER

    # L project, no materials in this branch.
    set_title(prs.slides[2], "L项目")
    set_body(prs.slides[2], "参与人员\n本期素材文件夹为空，暂无新增内容提交。", 24)
    set_body(prs.slides[3], "设计&沟通类：\n本期暂无提交。", 24)
    set_body(prs.slides[4], "资源制作类：\n本期暂无提交。", 24)

    # LL project.
    set_title(prs.slides[5], "重点工作1——LL项目")
    set_body(
        prs.slides[5],
        "参与人员\n音效设计 / 音频策划：关晗怡、安家辉\n音乐与PV相关：朱静雅、吉禹鹏\n音频技术：冯博宇、蔡新星",
        24,
    )
    set_title(prs.slides[6], "重点工作1——LL项目")
    set_body(
        prs.slides[6],
        "设计&沟通类：\n"
        "1. 河流画线功能优化需求：区域绘制后按 2*2 单元格自动拆分区块，减少手动画线成本。\n"
        "2. Manifest 冲突优化需求：减少多人传项目时拉新、重打包和重新生成 manifest 的重复操作。\n"
        "3. UI 编辑工具监听环境自动创建需求：进入场景后自动创建监听环境。\n"
        "4. 自动配置工具开发、Bug 修复与环境声画线相关功能重构。\n"
        "5. Jenkins 替换原有打包机流程方案验证，ClaudeCode 安装及使用。",
        18,
    )
    set_title(prs.slides[7], "重点工作1——LL项目")
    set_body(
        prs.slides[7],
        "资源制作类：\n"
        "1. 补充 UI / 地图交互物制作挂接，迭代部分资源与环境资源。\n"
        "2. 日配语音补充入版，语气词检查与重新标准化。\n"
        "3. 剧情资源优化：回忆滤镜增加语音混响，台词检查、资源剪辑、版本跑测与调优。\n"
        "4. 新版安德烈亚斯、马尔洛弹琴相关语音调试，解决重叠播放问题。\n"
        "5. PV 剪辑、Danger BGM Demo、PV 音乐剪辑版本等音乐 / 视频素材。",
        19,
    )
    ll_videos = [
        ROOT / "LL" / "安家辉" / "伊利萨尔-雄狮强袭.mp4",
        ROOT / "LL" / "安家辉" / "利昂战场技能.mp4",
        ROOT / "LL" / "安家辉" / "吉尔福德战场技能.mp4",
        ROOT / "LL" / "安家辉" / "埃拉战场技能吟唱.mp4",
        ROOT / "LL" / "安家辉" / "库莉娜大招.mp4",
        ROOT / "LL" / "安家辉" / "灯光秀.mp4",
        ROOT / "LL" / "安家辉" / "男主新战场技能.mp4",
        ROOT / "LL" / "朱静雅" / "Danger适配.mov",
    ]
    set_title(prs.slides[8], "重点工作1——LL项目｜产出展示 1")
    add_grid_videos(prs.slides[8], ll_videos[:4], [p.stem for p in ll_videos[:4]])
    set_title(prs.slides[9], "重点工作1——LL项目｜产出展示 2")
    add_grid_videos(prs.slides[9], ll_videos[4:8], [p.stem for p in ll_videos[4:8]])
    set_title(prs.slides[10], "重点工作1——LL项目｜产出展示 3")
    clear_body(prs.slides[10])
    add_video(
        prs.slides[10],
        ROOT / "LL" / "吉禹鹏" / "PV剪辑.mp4",
        Inches(0.7),
        Inches(1.2),
        Inches(7.3),
        Inches(4.35),
        "PV剪辑",
    )
    add_audio_list(
        prs.slides[10],
        "音频素材",
        ["Discussion.mp3", "0513_PV音乐剪辑版本_BC同一首02.mp3", "0515DEMO_BGM_Story_Danger.wav"],
        Inches(8.35),
        Inches(1.35),
        Inches(4.1),
        Inches(3.6),
    )

    # U project, shifted by two slides after LL insertions.
    set_title(prs.slides[11], "重点工作2——U项目")
    set_body(
        prs.slides[11],
        "参与人员：\n音效师 / 音频策划：郑青\n音频技术：冯博宇、蔡新星",
        24,
    )
    set_body(
        prs.slides[12],
        "资源制作类：\n"
        "1. 音效 40 个。\n"
        "2. 语音 47 句。\n"
        "3. 视频 2 条。\n"
        "4. 语音发包回包。\n"
        "5. Wwise 维护。\n"
        "6. Unity 挂接。",
        22,
    )
    set_title(prs.slides[13], "重点工作2——U项目｜产出展示")
    clear_body(prs.slides[13])
    add_video(prs.slides[13], ROOT / "U" / "郑青" / "试听.mp4", Inches(1.2), Inches(1.25), Inches(10.8), Inches(5.25), "试听")

    # FM project.
    set_title(prs.slides[14], "重点工作3——FM项目")
    set_body(
        prs.slides[14],
        "参与人员\n音效设计 / 音乐设计 / 音频策划：冯博宇\n音频技术：冯博宇、蔡新星",
        24,
    )
    set_body(prs.slides[15], "设计&沟通类：\nFM：EN9 线上 Bug、OB23 音效制作、挂接与测试。", 24)
    set_body(prs.slides[16], "资源制作类：\nUI：8 个", 26)
    set_body(prs.slides[17], "产出展示类：", 22)
    add_video(prs.slides[17], ROOT / "FM" / "冯博宇.mp4", Inches(1.15), Inches(1.3), Inches(10.6), Inches(5.1), "FM UI 音效")

    # EF project, no materials.
    set_title(prs.slides[18], "EF项目")
    set_body(prs.slides[18], "参与人员\n本期素材文件夹为空，暂无新增内容提交。", 24)
    set_title(prs.slides[19], "EF项目")
    set_body(prs.slides[19], "设计&沟通类：\n本期暂无提交。", 24)
    set_title(prs.slides[20], "EF项目")
    set_body(prs.slides[20], "资源制作类：\n本期暂无提交。", 24)

    # TA / audio tech.
    set_title(prs.slides[21], "音频技术")
    set_body(
        prs.slides[21],
        "平台：工具维护。\n"
        "FM：EN9 线上 Bug、OB23 音效制作、挂接与测试。\n"
        "U：TW OB64 版本工单。\n"
        "LL：自动配置工具开发、Bug 修复，环境声画线相关功能重构。\n"
        "平台流程：尝试通过 Jenkins 替换原有打包机流程。\n"
        "工具学习：ClaudeCode 安装及使用。",
        20,
    )

    # Other and thanks.
    set_title(prs.slides[22], "其他内容")
    set_body(prs.slides[22], "团队建设 / 其他事项：\n本期暂无额外提交。", 24)
    thanks_shape = content_shape(prs.slides[23]) or title_shape(prs.slides[23])
    set_text(thanks_shape, "THANK YOU", font_size=38, bold_first=True)
    for p in thanks_shape.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.color.rgb = RGBColor(31, 78, 121)

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(OUT_PPT)
    print("slides", len(prs.slides))


if __name__ == "__main__":
    main()
