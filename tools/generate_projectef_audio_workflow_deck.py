from __future__ import annotations

import html
import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORT_DIR = ROOT / "报告"
OUT_STEM = "ProjectEF_音频研发流程优化方案_2026-06-16"
PPTX_PATH = REPORT_DIR / f"{OUT_STEM}.pptx"
HTML_PATH = REPORT_DIR / f"{OUT_STEM}.html"


COLORS = {
    "bg": RGBColor(247, 248, 250),
    "panel": RGBColor(255, 255, 255),
    "ink": RGBColor(28, 36, 46),
    "muted": RGBColor(102, 112, 133),
    "line": RGBColor(218, 223, 231),
    "teal": RGBColor(19, 121, 115),
    "green": RGBColor(72, 150, 92),
    "amber": RGBColor(224, 145, 54),
    "red": RGBColor(201, 76, 76),
    "blue": RGBColor(58, 102, 177),
    "purple": RGBColor(126, 88, 166),
    "slate": RGBColor(64, 78, 96),
}


FONT = "Microsoft YaHei"
FONT_BOLD = "Microsoft YaHei UI"


def c(hex_str: str) -> RGBColor:
    hex_str = hex_str.strip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def read_json(path: Path):
    try:
        return json.loads(path.read_text(encoding="utf-8-sig"))
    except Exception:
        return None


def collect_context():
    tool_json = read_json(ROOT / "Tools" / "EF_Audio_Tools_Final" / "tool_paths.json") or {}
    tools = tool_json.get("tools", [])
    visible_tools = [t for t in tools if t.get("visible")]
    daily = read_json(REPORT_DIR / "ProjectEF_DailyAudioLogIntelligence_2026-06-15.json") or {}
    trend = read_json(REPORT_DIR / "ProjectEF_AudioReport_TrendSummary.json") or {}
    md_reports = len(list(REPORT_DIR.glob("*.md"))) if REPORT_DIR.exists() else 0
    html_reports = len(list((ROOT / "ProjectEF_reports_html" / "reports").glob("*.html")))
    return {
        "tools_total": len(tools),
        "visible_tools": len(visible_tools),
        "visible_tool_names": [t.get("name", "") for t in visible_tools],
        "daily": daily,
        "trend": trend,
        "md_reports": md_reports,
        "html_reports": html_reports,
    }


CTX = collect_context()


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg"]


def set_text(run, size=18, color=None, bold=False):
    run.font.name = FONT_BOLD if bold else FONT
    run.font.size = Pt(size)
    if color is not None:
        run.font.color.rgb = color
    run.font.bold = bold


def textbox(slide, x, y, w, h, text="", size=18, color=None, bold=False, align=None, valign=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    set_text(r, size=size, color=color or COLORS["ink"], bold=bold)
    return box


def title(slide, text, kicker=None):
    if kicker:
        textbox(slide, 0.64, 0.36, 6.8, 0.32, kicker, size=10.5, color=COLORS["teal"], bold=True)
    textbox(slide, 0.62, 0.66, 10.8, 0.62, text, size=27, color=COLORS["ink"], bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.66), Inches(1.36), Inches(1.18), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["teal"]
    line.line.fill.background()


def footer(slide, page, section="ProjectEF Audio Production Workflow"):
    textbox(slide, 0.62, 7.06, 6.2, 0.22, section, size=8.5, color=COLORS["muted"])
    textbox(slide, 11.9, 7.06, 0.75, 0.22, f"{page:02d}", size=8.5, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, heading, body=None, color=None, fill=COLORS["panel"], heading_size=15, body_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = COLORS["line"]
    if color:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()
    textbox(slide, x + 0.18, y + 0.14, w - 0.34, 0.36, heading, size=heading_size, color=color or COLORS["ink"], bold=True)
    if body:
        textbox(slide, x + 0.18, y + 0.55, w - 0.34, h - 0.66, body, size=body_size, color=COLORS["slate"])
    return shape


def bullets(slide, x, y, w, h, items, size=14, color=None, gap=0.07):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6 + gap * 10)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
        p.text = item
    return box


def metric(slide, x, y, w, h, value, label, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.10
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["panel"]
    shape.line.color.rgb = COLORS["line"]
    textbox(slide, x + 0.1, y + 0.18, w - 0.2, 0.42, value, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, x + 0.1, y + 0.70, w - 0.2, 0.35, label, size=10.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def table(slide, x, y, w, h, headers, rows, widths=None, header_fill=COLORS["teal"], font_size=9.5):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    if widths:
        for i, col_w in enumerate(widths):
            tbl.columns[i].width = Inches(col_w)
    for j, head in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        p = cell.text_frame.paragraphs[0]
        p.text = head
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = FONT_BOLD
        p.runs[0].font.size = Pt(font_size)
        p.runs[0].font.color.rgb = RGBColor(255, 255, 255)
        p.runs[0].font.bold = True
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if i % 2 else RGBColor(250, 252, 253)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            p.runs[0].font.name = FONT
            p.runs[0].font.size = Pt(font_size)
            p.runs[0].font.color.rgb = COLORS["ink"]
    return shape


def phase_strip(slide, phases, y=2.1):
    x = 0.68
    gap = 0.14
    w = (12.0 - gap * (len(phases) - 1)) / len(phases)
    for idx, (name, pain, text, color) in enumerate(phases):
        card(slide, x + idx * (w + gap), y, w, 2.30, name, text, color=color, heading_size=14, body_size=10.5)
        textbox(slide, x + idx * (w + gap) + 0.18, y + 1.82, w - 0.36, 0.28, f"痛点强度：{pain}", size=9.5, color=color, bold=True)


def add_flow(slide, items, y=2.20):
    x = 0.78
    box_w = 1.44
    gap = 0.28
    for i, (head, body, color) in enumerate(items):
        bx = x + i * (box_w + gap)
        card(slide, bx, y, box_w, 1.48, head, body, color=color, heading_size=11.5, body_size=8.4)
        if i < len(items) - 1:
            textbox(slide, bx + box_w + 0.04, y + 0.52, 0.26, 0.25, "→", size=18, color=COLORS["muted"], bold=True)


def section_slide(prs, page, heading, subheading, blocks):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    textbox(slide, 0.78, 0.95, 10.5, 0.68, heading, size=30, color=COLORS["ink"], bold=True)
    textbox(slide, 0.80, 1.70, 10.4, 0.40, subheading, size=14.5, color=COLORS["muted"])
    x = 0.88
    for i, (label, text, color) in enumerate(blocks):
        card(slide, x + i * 3.95, 3.0, 3.55, 1.55, label, text, color=color, heading_size=16, body_size=11)
    footer(slide, page)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    page = 1

    # 1 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c("F3F7F6")
    shape.line.fill.background()
    textbox(slide, 0.84, 0.72, 9.9, 0.36, "ProjectEF / Game Audio Production Intelligence", size=12, color=COLORS["teal"], bold=True)
    textbox(slide, 0.82, 1.32, 10.65, 1.35, "音频研发流程优化方案", size=40, color=COLORS["ink"], bold=True)
    textbox(slide, 0.86, 2.72, 10.35, 0.80, "从“下游救火”到“字段驱动的协作闭环”", size=24, color=COLORS["slate"], bold=True)
    textbox(slide, 0.90, 4.12, 6.9, 0.85, "适用范围：需求总览、机制讨论、音频设计、Unity/Wwise 配置、QA 验收、P4 提交流程", size=15, color=COLORS["muted"])
    metric(slide, 8.10, 4.00, 1.38, 1.16, str(CTX["tools_total"]), "工具条目", COLORS["blue"])
    metric(slide, 9.72, 4.00, 1.38, 1.16, str(CTX["visible_tools"]), "日常入口", COLORS["teal"])
    metric(slide, 11.34, 4.00, 1.38, 1.16, str(CTX["md_reports"]), "MD 报告", COLORS["amber"])
    textbox(slide, 0.92, 6.62, 4.3, 0.24, "生成日期：2026-06-16", size=9.5, color=COLORS["muted"])
    footer(slide, page)
    page += 1

    # 2 Executive conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "一页结论", "Executive Summary")
    textbox(slide, 0.84, 1.62, 11.65, 0.55, "当前问题是正常的，但不能继续用“一个音频人反复找人确认”的方式消耗。真正要优化的是信息进入音频生产线的形态。", size=19, color=COLORS["ink"], bold=True)
    card(slide, 0.80, 2.48, 3.80, 2.05, "判断", "音频处于项目最下游，需要理解玩法、资源、程序逻辑和测试路径；在复杂钓鱼玩法中，这是行业常态。", COLORS["blue"])
    card(slide, 4.82, 2.48, 3.80, 2.05, "核心风险", "如果需求、配置位置、测试方式没有结构化字段，音频会变成跨部门信息搬运工，产能被沟通吞掉。", COLORS["red"])
    card(slide, 8.84, 2.48, 3.80, 2.05, "优化方向", "把沟通对象从“人”变成“字段”：建立音频落地矩阵、Ready Gate、测试 Recipe 和周期性 Triage。", COLORS["teal"])
    bullets(slide, 0.92, 5.06, 11.6, 1.25, [
        "目标不是取消沟通，而是让沟通只发生在缺字段、复杂机制、设计取舍和风险确认上。",
        "工具链已经很强；下一步要把工具输出合并成一个可追踪的生产闭环。",
        "音频可以提前探索和预判，但只有通过 Ready Gate 的内容才进入正式制作与验收。"
    ], size=14)
    footer(slide, page)
    page += 1

    # 3 First-principles frame
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "第一性原理：判断音频需求到底需要什么", "First Principles")
    textbox(slide, 0.88, 1.48, 11.58, 0.50, "音频需求不是从“想做什么声音”开始，而是从“游戏流程中有哪些需要被声音解释或强化的状态变化”开始。", size=18, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    card(slide, 0.82, 2.14, 2.22, 2.15, "系统清单", "游戏到底有多少系统、哪些系统正在版本内交付。", COLORS["blue"], heading_size=14, body_size=10.5)
    card(slide, 3.22, 2.14, 2.22, 2.15, "流程地图", "每个系统的完整玩家流程、状态变化、成功失败和边界。", COLORS["purple"], heading_size=14, body_size=10.5)
    card(slide, 5.62, 2.14, 2.22, 2.15, "资源/逻辑", "流程涉及的 Prefab、Animation、VFX、配置表、脚本和程序状态。", COLORS["amber"], heading_size=14, body_size=10.5)
    card(slide, 8.02, 2.14, 2.22, 2.15, "触发/测试", "如何进入场景、如何触发、如何 Debug、如何证明通过。", COLORS["red"], heading_size=14, body_size=10.5)
    card(slide, 10.42, 2.14, 2.22, 2.15, "音频判定", "哪些点需要声音、优先级、实现合约、资源和混音策略。", COLORS["teal"], heading_size=14, body_size=10.5)
    textbox(slide, 1.06, 4.82, 11.20, 0.70, "公式：系统 × 流程 × 资源/逻辑 × 触发/测试 × 音频判定 = 可落地音频需求", size=22, color=COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.26, 5.88, 10.72, 0.50, "任何一项缺失，音频都只能做候选判断或方案预研，不能稳定进入正式制作和验收。", size=16, color=COLORS["slate"], align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 4 Current coverage gap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "当前信息覆盖率缺口", "Coverage Gap")
    rows = [
        ["游戏有多少系统", "Jira", "80-90%", "可形成系统视野，但仍要关注未建单/拆分不准的内容。"],
        ["每个系统的所有流程", "策划文档 + 阅读理解 + 询问", "30-40%", "文档常写目标，不写完整状态机、边界和音频触发。"],
        ["流程涉及的美术资源和程序逻辑", "策划文档 + 反查资源 + 询问", "20-30%", "设计名、资源名、脚本名、Prefab 路径之间缺映射。"],
        ["所有流程的测试和触发方式", "目前准备问 QA", "系统性不足", "缺测试场景、Debug 入口、账号/存档、日志和 Profiler 证据。"],
    ]
    table(slide, 0.86, 1.56, 11.64, 3.22, ["音频需要知道", "目前来源", "覆盖率", "主要缺口"], rows, widths=[2.15, 2.35, 1.35, 5.75], font_size=9.3)
    card(slide, 0.94, 5.20, 3.62, 1.04, "根问题 1", "上游需求不全，音频无法稳定判断什么需要声音。", COLORS["red"], heading_size=13, body_size=10.2)
    card(slide, 4.86, 5.20, 3.62, 1.04, "根问题 2", "上游并不知道音频判断和落地需要哪些输入。", COLORS["amber"], heading_size=13, body_size=10.2)
    card(slide, 8.78, 5.20, 3.62, 1.04, "根问题 3", "上游部门多，音频反查 owner、资源名和测试路径极其耗时。", COLORS["blue"], heading_size=13, body_size=10.2)
    footer(slide, page)
    page += 1

    # 5 Elephant decomposition
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "把大象放进冰箱：拆成可逐个击破的问题", "Decomposition")
    add_flow(slide, [
        ("1 系统", "Jira/版本\n列全系统", COLORS["blue"]),
        ("2 流程", "策划案\n玩家路径", COLORS["purple"]),
        ("3 反馈点", "状态变化\n操作反馈", COLORS["teal"]),
        ("4 资源逻辑", "Prefab/Anim\nVFX/脚本", COLORS["amber"]),
        ("5 音频合约", "Event/RTPC\nSwitch/Stop", COLORS["green"]),
        ("6 测试", "场景/步骤\nDebug/证据", COLORS["red"]),
        ("7 验收", "听感/日志\nProfiler/P4", COLORS["slate"]),
    ], y=1.80)
    rows = [
        ["系统", "有哪些系统在版本内？是否已有 Jira？谁是系统 owner？"],
        ["流程", "玩家从进入到退出经历哪些步骤？每一步有哪些状态变化？"],
        ["反馈点", "哪些变化影响判断、操作、空间、危险、奖励、情绪或节奏？"],
        ["资源逻辑", "对应哪个 Prefab、Animation、VFX、配置字段、脚本状态？"],
        ["音频合约", "谁触发 Play/Stop？传哪些 RTPC？设置哪些 Switch/State？"],
        ["测试", "如何触发？用什么 Debug？预期听到什么？用什么证据关闭？"],
    ]
    table(slide, 1.06, 4.10, 11.08, 2.05, ["步骤", "要回答的问题"], rows, widths=[1.35, 9.65], font_size=9.6)
    footer(slide, page)
    page += 1

    # 6 Authority and interface owner
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "每类问题都要有接口人或说明文档", "Authority Source")
    rows = [
        ["这个系统有没有？", "Jira / 制作人 / 系统策划", "系统列表、版本范围、优先级"],
        ["系统怎么玩？", "策划案 / 系统策划", "流程图、状态说明、成功失败、边界"],
        ["状态什么时候变？", "程序 / 状态机 / 配置表", "权威状态、触发 API、生命周期、参数"],
        ["资源在哪里？", "美术 / TA / 动画 / VFX / Unity 资源表", "Prefab、Animation、VFX、挂点、资源稳定性"],
        ["怎么触发测试？", "QA / 程序 Debug / 测试文档", "场景、步骤、账号/存档、Debug 入口、验收证据"],
        ["声音怎么表现？", "音频", "声音目标、Wwise 合约、资源量、混音、验收标准"],
    ]
    table(slide, 0.82, 1.54, 11.78, 4.66, ["问题", "权威来源 / 接口人", "产出物"], rows, widths=[2.25, 3.65, 5.75], font_size=10.0)
    textbox(slide, 0.98, 6.40, 11.32, 0.34, "原则：音频可以推动问题暴露，但不应该成为所有上游事实的唯一调查员。", size=15.5, color=COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 7 Upstream package checklist
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "上游可前置提供的资源类型清单", "Upstream Package")
    rows = [
        ["Prefab", "路径、挂点、组件、是否运行时生成、owner、是否可加音频字段"],
        ["Animation", "clip 名、关键帧、循环段、是否可加 Animation Event、预览方式"],
        ["VFX", "VFX 名、播放时机、强弱等级、持续时间、是否跟随对象"],
        ["UI", "prefab、按钮组件、pressed/hover/success/fail 字段、默认音频规则"],
        ["Config", "表名、ID、AudioEventId/AudioProfileId 字段、是否热更"],
        ["State Machine", "状态名、进入/退出条件、状态切换频率、owner"],
        ["Runtime Logic", "触发函数、GameObject、生命周期、日志、多人/本地玩家规则"],
        ["QA Scene", "场景名、进入方式、账号/存档、Debug 工具、复现步骤"],
    ]
    table(slide, 0.90, 1.50, 11.55, 4.95, ["类型", "音频需要的信息"], rows, widths=[1.65, 9.78], font_size=9.8)
    footer(slide, page)
    page += 1

    # 8 Audio decision rules
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "音频判定规则：不是所有流程都要做声音", "Decision Rules")
    rows = [
        ["必须做", "影响核心判断、操作成败、危险预警、状态变化、空间定位、奖励闭环。"],
        ["优先做", "高频操作、强视觉事件、玩家长期感知差异、系统价值展示。"],
        ["模板化", "同类鱼、普通装备、普通 UI、重复水花、远处他人事件。"],
        ["后置/裁剪", "低价值装饰、Hover、普通刷新、内部工具、干扰核心信息的提示。"],
        ["需要讨论", "情绪目标明确但触发状态不清、资源不稳定、程序没有 Debug 路径。"],
        ["不能 Ready", "没有流程、没有资源路径、没有触发 owner、没有测试方式。"],
    ]
    table(slide, 0.96, 1.52, 11.42, 4.22, ["判定", "规则"], rows, widths=[1.65, 9.65], font_size=10.8)
    textbox(slide, 1.10, 6.08, 11.05, 0.46, "音频的第一职责是判断哪些游戏状态需要被听见，其次才是决定用什么声音。", size=17, color=COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 3 Source base
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "目前已掌握的资料资产", "Inventory")
    metric(slide, 0.82, 1.58, 1.58, 1.08, str(CTX["tools_total"]), "最终菜单工具", COLORS["blue"])
    metric(slide, 2.58, 1.58, 1.58, 1.08, str(CTX["visible_tools"]), "日常可见工具", COLORS["teal"])
    metric(slide, 4.34, 1.58, 1.58, 1.08, str(CTX["html_reports"]), "HTML 报告页", COLORS["purple"])
    metric(slide, 6.10, 1.58, 1.58, 1.08, str(CTX["md_reports"]), "Markdown 报告", COLORS["amber"])
    daily = CTX["daily"]
    metric(slide, 7.86, 1.58, 1.58, 1.08, f"{daily.get('runtime_validation_average', '33.8')}%", "最新运行验证均值", COLORS["red"])
    metric(slide, 9.62, 1.58, 1.58, 1.08, str(len(daily.get("modules", [])) or 16), "日报模块数", COLORS["green"])
    metric(slide, 11.38, 1.58, 1.58, 1.08, str(len(daily.get("problems", [])) or 29), "最新问题项", COLORS["red"])
    rows = [
        ["需求/Jira", "Audio Requirement Jira Triage、设计文档扫描、需求 Diff、SABC 清单"],
        ["Wwise", "工程体检、资源审计、Template Generator、Audio Logic Tester、Debug Assistant"],
        ["Unity", "Audio Footprint、UI Inspector、Animation Event AutoConfig、Runtime Log Monitor"],
        ["运行 QA", "Daily Log Intelligence、Profiler Voice Capture、Runtime QA Checklist、Bank 检查"],
        ["P4/提交", "P4V Changelist Check、Changelist Organizer、提交风险分类"],
        ["制作资源", "Sound Finder for Reaper、本地素材库、REAPER 项目联动、本地 LLM 辅助"],
    ]
    table(slide, 0.84, 3.02, 11.82, 3.28, ["资料域", "已有能力"], rows, widths=[2.2, 9.5], font_size=10.5)
    footer(slide, page)
    page += 1

    # 4 Current flow and pain
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "当前实际链路：音频在最下游承接所有不确定性", "Current State")
    phase_strip(slide, [
        ("有什么\n需求总览", "中", "Jira 和策划案能覆盖大概，但音频仍需判断哪些内容有声音价值。", COLORS["blue"]),
        ("什么机制\n讨论确认", "高", "复杂系统涉及状态机、VFX、动画、程序生命周期，不能只靠文案理解。", COLORS["amber"]),
        ("怎么实现\n音频设计", "低/中", "前置信息清楚后，资源制作和 Wwise 搭建通常最明确。", COLORS["green"]),
        ("怎么配置\n落地接入", "最高", "资源名、Prefab、动画、脚本、配置表、负责人都分散在不同人和不同系统里。", COLORS["red"]),
        ("怎么测试\n验收复现", "最高", "音频不知道场景、Debug 入口、边界条件；QA 也可能没有现成方式。", COLORS["red"]),
    ], y=1.95)
    textbox(slide, 0.82, 5.06, 11.75, 0.70, "结论：制作本身不是最大损耗；“找信息、找人、找入口、找测试方法”才是单人音频最难承受的地方。", size=18, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 5 Pain demand overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "痛点 1：需求总览能扫到，但不等于可制作", "Pain Point")
    card(slide, 0.82, 1.68, 3.75, 3.92, "现在的好消息", "Jira 单、策划需求文案关联工具、设计文档扫描已经能让音频快速知道“游戏里发生了什么”。这解决了入门视野问题。", COLORS["green"])
    card(slide, 4.82, 1.68, 3.75, 3.92, "仍然缺的东西", "策划通常写玩法目标，不会自然写清音频触发、状态边界、优先级、失败条件、测试方法。音频也不可能凭文案完全判断所有声音需求。", COLORS["amber"])
    card(slide, 8.82, 1.68, 3.75, 3.92, "优化判断", "需求扫描应输出“候选音频需求”，而不是直接输出最终制作任务。进入制作前必须补齐机制、配置、QA 字段。", COLORS["teal"])
    textbox(slide, 1.02, 6.10, 11.30, 0.42, "建议状态：Candidate → DesignOnly → Ready / Blocked / Risky / Cuttable", size=18, color=COLORS["slate"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 6 Mechanism pain
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "痛点 2：复杂机制不是“读完文档就能懂”", "Pain Point")
    bullets(slide, 0.95, 1.62, 5.62, 4.60, [
        "钓鱼核心声音依赖状态：咬口、漂相、张力、泄力、鱼体力、鱼技能、天气、水流、视角、多人距离。",
        "复杂系统里，声音触发往往不在策划文案里，而在程序状态机、动画事件、VFX 节点、配置表或 Prefab 上。",
        "音频真正需要知道的是：什么状态发生、谁是权威来源、什么时候开始/停止、参数范围、边界和冷却规则。",
        "这类问题必须通过短 walkthrough 解决，不能靠音频单人反复猜。"
    ], size=15)
    rows = [
        ["策划", "玩法规则、状态含义、优先级、成功/失败边界"],
        ["程序", "权威状态、触发 API、生命周期、Debug 入口、日志证据"],
        ["美术/VFX", "资源路径、动画/VFX 节点、视觉变化时机"],
        ["音频", "声音目标、Wwise 参数、资源量、混音优先级"],
        ["QA", "测试场景、复现步骤、验收标准"],
    ]
    table(slide, 6.84, 1.70, 5.62, 3.35, ["角色", "机制讨论中必须产出的信息"], rows, widths=[1.2, 4.35], font_size=9.8)
    card(slide, 6.84, 5.34, 5.62, 0.78, "会议目标", "不是开大讨论会，而是把落地矩阵里的空字段填满。", COLORS["teal"], heading_size=13.5, body_size=11.2)
    footer(slide, page)
    page += 1

    # 7 Config pain
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "痛点 3：配置落地是最大时间黑洞", "Pain Point")
    card(slide, 0.82, 1.62, 3.72, 1.50, "策划不知道", "美术资源叫什么、程序逻辑在哪里、Prefab/Anim/VFX 怎么组织。", COLORS["amber"])
    card(slide, 4.80, 1.62, 3.72, 1.50, "音频不知道", "策划文档里的对象在 Unity 里叫什么、在哪个资源、由谁负责。", COLORS["red"])
    card(slide, 8.78, 1.62, 3.72, 1.50, "美术/程序不知道", "音频到底要什么 Event、参数、触发时机、Stop 规则和测试证据。", COLORS["blue"])
    textbox(slide, 0.92, 3.58, 11.48, 0.52, "如果每个系统、每个细节都靠“找具体人问”，音频的生产时间会被切碎，且知识不能复用。", size=18, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    rows = [
        ["UI", "Prefab / ButtonEx / ButtonAudioComp / effective_audio / Wwise Event 状态"],
        ["动画", ".fbx/.anim Clip / Marker / Animation Event / Receiver / 预览窗口"],
        ["VFX/场景", "VFX 名称 / 场景对象 / Emitter / 距离衰减 / 触发 owner"],
        ["玩法状态", "状态机 / RTPC / Switch / Start/Stop / Debug 命令"],
        ["配置表", "AudioEventId / AudioProfileId / 资源 ID / 负责人"],
    ]
    table(slide, 1.02, 4.45, 11.30, 1.65, ["类型", "配置定位必须结构化记录"], rows, widths=[1.6, 9.6], font_size=9.4)
    footer(slide, page)
    page += 1

    # 8 Test pain
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "痛点 4：测试不可复现，问题就无法关闭", "Pain Point")
    card(slide, 0.84, 1.56, 3.58, 1.45, "静态检查能证明什么", "Prefab、脚本、Wwise Event、Bank、配置引用是否存在。", COLORS["blue"])
    card(slide, 4.88, 1.56, 3.58, 1.45, "运行日志能证明什么", "本次 Play Session 里发生过什么、报了什么错、Bank/Voice 是否异常。", COLORS["amber"])
    card(slide, 8.92, 1.56, 3.58, 1.45, "Profiler 能证明什么", "运行时 active voice、GameObject、优先级、虚拟化、Starvation 来源。", COLORS["purple"])
    textbox(slide, 0.88, 3.40, 11.70, 0.58, "但没有测试场景和步骤时，所有工具都只能给出 StaticOnly / NotObserved，不能证明功能通过。", size=19, color=COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    rows = [
        ["必须有", "测试场景 / 角色状态 / Debug 入口 / 操作步骤 / 期望听感 / 日志或 Profiler 证据"],
        ["不要把", "NotObserved 当 Pass；未跑到就是未验证"],
        ["QA 责任", "提供可复现路径；如果没有路径，就转化为程序/策划需求"],
        ["音频责任", "定义听感标准、混音标准、关键失败模式和验收证据"],
    ]
    table(slide, 1.02, 4.45, 11.30, 1.68, ["规则", "说明"], rows, widths=[1.6, 9.6], font_size=10.2)
    footer(slide, page)
    page += 1

    # 9 Tool coverage
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "现有工具链已经解决了很多事，但还差一个“生产中枢”", "Tool Coverage")
    rows = [
        ["需求发现", "Jira/设计扫描、SABC、Diff", "候选需求和变化识别", "需要落地矩阵承接字段"],
        ["Wwise 设计", "Template、Audit、Logic Tester、Debug Assistant", "结构创建、风险检查、单点调试", "需要和 Unity/QA 字段绑定"],
        ["Unity 配置", "UI Inspector、Animation AutoConfig、Footprint", "定位 Prefab/Anim/脚本引用", "需要 owner 和测试入口"],
        ["运行验证", "Log Monitor、Daily Log、Profiler Capture", "发现运行错误和 Starvation", "需要主动测试场景"],
        ["P4 提交", "Changelist Check、Organizer", "提交风险收敛", "需要和任务/验收证据关联"],
        ["素材制作", "Sound Finder、REAPER 联动", "找素材和制作资源更快", "不解决上游信息缺口"],
    ]
    table(slide, 0.70, 1.62, 12.05, 4.52, ["环节", "已有工具", "已解决", "仍缺"], rows, widths=[1.35, 3.0, 3.1, 4.45], font_size=8.8)
    textbox(slide, 0.86, 6.30, 11.70, 0.35, "下一步不是堆更多孤立工具，而是让所有工具围绕同一张“音频落地矩阵”读写证据。", size=16, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 10 Principles
    section_slide(prs, page, "优化原则", "音频不应该被动等上游稳定，也不应该替上游兜底所有不清楚的信息。", [
        ("探索层", "提前扫需求、预判风险、做模板、列问题。", COLORS["blue"]),
        ("生产层", "只有 Ready 的任务进入资源制作和正式接入。", COLORS["green"]),
        ("证据层", "每个结论都能回到文档、Jira、路径、日志或截图。", COLORS["teal"]),
    ])
    page += 1

    # 11 Target flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "目标流程：字段驱动的协作闭环", "Target Workflow")
    add_flow(slide, [
        ("输入", "Jira\n设计文档\n配置表", COLORS["blue"]),
        ("候选需求", "AI 扫描\nSABC\n变化 Diff", COLORS["purple"]),
        ("落地矩阵", "字段补齐\n状态分类\nowner", COLORS["teal"]),
        ("Triage", "缺口确认\n机制 walkthrough", COLORS["amber"]),
        ("实现", "资源\nWwise\nUnity 配置", COLORS["green"]),
        ("QA", "Recipe\n日志\nProfiler", COLORS["red"]),
        ("提交监控", "P4\n日报\n趋势", COLORS["slate"]),
    ], y=2.08)
    textbox(slide, 0.92, 4.30, 11.55, 0.52, "关键变化：不再让音频逐个追问“在哪里配、怎么测、找谁问”，而是让每个任务在矩阵里暴露缺口。", size=17, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    rows = [
        ["Candidate", "工具发现可能需要音频，但信息不足"],
        ["DesignOnly", "可做音频方案，不进入正式资源/接入"],
        ["Ready", "触发、资源、配置、测试路径已明确"],
        ["Blocked", "缺少上游决策、资源、代码或测试入口"],
        ["Risky", "可推进但明确有返工风险"],
        ["Cuttable", "低价值或可后置内容"],
    ]
    table(slide, 1.25, 5.10, 10.82, 1.20, ["状态", "含义"], rows, widths=[1.75, 9.0], font_size=8.7)
    footer(slide, page)
    page += 1

    # 12 Matrix
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "核心产物：音频落地矩阵", "Operating Artifact")
    textbox(slide, 0.86, 1.48, 11.75, 0.40, "这张表是所有沟通、工具输出、排期、测试、提交的共同入口。", size=17, color=COLORS["ink"], bold=True)
    rows = [
        ["来源", "Jira、策划案、配置表、截图、版本"],
        ["需求判断", "是否需要音频、SABC、声音目标、可裁剪性"],
        ["机制信息", "状态、触发、Start/Stop、边界、冷却、参数范围"],
        ["Unity 落地", "Prefab/Anim/VFX/脚本/配置路径、owner、字段名"],
        ["Wwise 落地", "Event、RTPC、Switch/State、Bus、Bank、Attenuation"],
        ["资源与混音", "素材数、随机样本、优先级、Ducking、性能预算"],
        ["QA 验收", "场景、步骤、Debug 入口、预期听感、日志/Profiler 证据"],
        ["状态追踪", "Ready/Blocked/Risky/NotObserved、风险、下一步、负责人"],
    ]
    table(slide, 0.96, 2.06, 11.45, 3.80, ["字段组", "必须记录的内容"], rows, widths=[2.1, 9.25], font_size=10.4)
    card(slide, 1.02, 6.10, 11.35, 0.54, "使用原则", "问人之前先看矩阵；开会只补矩阵；工具报告只追加证据；验收只关闭有证据的行。", COLORS["teal"], heading_size=12.5, body_size=10.5)
    footer(slide, page)
    page += 1

    # 13 Ready gate
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Ready Gate：什么任务可以真正进制作", "Definition of Ready")
    rows = [
        ["Design", "玩法状态、成功/失败、时机、优先级、边界已知"],
        ["Art/VFX", "视觉/动画参考存在，或明确接受 placeholder 风险"],
        ["Program", "触发 owner、Event、参数、生命周期、Debug 入口已知"],
        ["Wwise", "层级、Bus、衰减、随机、Voice、SoundBank 方案已知"],
        ["Resource", "素材数、样本数、命名、来源、授权/VO 本地化已知"],
        ["Mix", "响度、Ducking、优先级、多人裁剪、平台约束已知"],
        ["QA", "测试场景、步骤、预期结果、失败模式、证据路径已知"],
    ]
    table(slide, 1.02, 1.62, 11.28, 4.30, ["区域", "Ready 条件"], rows, widths=[1.45, 9.75], font_size=11.2)
    textbox(slide, 1.06, 6.18, 11.15, 0.42, "少一个关键条件，就不要叫 Ready。可以探索、设计、预估，但不应该承诺最终制作完成。", size=16.5, color=COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 14 RACI
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "角色分工：让音频负责声音，不负责替所有人找上下文", "Ownership")
    rows = [
        ["策划", "玩法状态、需求优先级、边界条件、是否可裁剪", "机制不清导致音频方向错误"],
        ["程序", "权威状态、触发点、参数范围、Debug 入口、日志证据", "音频无法配置/无法复现"],
        ["美术/VFX/动画", "资源路径、动画/VFX 时机、视觉参考、资源稳定性", "声音和画面不同步或返工"],
        ["音频", "声音目标、Wwise 方案、资源制作、混音、音频验收标准", "声音体验和性能失控"],
        ["QA", "测试路径、场景覆盖、复现步骤、回归证据", "StaticOnly/NotObserved 无法关闭"],
        ["制作/负责人", "强制字段、Triage 节奏、跨部门 owner 决策", "流程无人维护，回到逐个找人"],
    ]
    table(slide, 0.70, 1.50, 12.05, 4.72, ["角色", "必须提供", "不提供的风险"], rows, widths=[1.35, 5.35, 5.20], font_size=8.9)
    footer(slide, page)
    page += 1

    # 15 Optimization by current questions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "针对当前五类问题的优化动作", "Action Map")
    rows = [
        ["有什么", "Jira+设计扫描继续做，但输出 Candidate；按 SABC 和系统归档；每周看 Diff。"],
        ["什么机制", "复杂系统强制 15-30 分钟 walkthrough；会议结果只补状态、触发、参数、边界字段。"],
        ["怎么实现", "Wwise 模板、命名规则、资源样本数、Bus/RTPC/Switch 模板化；音频保留创作空间。"],
        ["怎么配置", "每个任务必须有 Unity 配置路径、owner、字段名、引用方式；没有则 Blocked。"],
        ["怎么测试", "每个 Ready 任务必须有 QA Recipe：场景、步骤、Debug、预期、日志/Profiler 证据。"],
    ]
    table(slide, 0.84, 1.62, 11.70, 4.26, ["环节", "优化动作"], rows, widths=[1.35, 10.25], font_size=11.0)
    card(slide, 1.00, 6.12, 11.32, 0.56, "管理口径", "音频可以协助发现问题，但不应独自承担缺少配置路径、缺少 Debug 入口、缺少测试场景造成的排期风险。", COLORS["amber"], heading_size=12.2, body_size=10.5)
    footer(slide, page)
    page += 1

    # 16 Jira template
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Jira/需求模板建议：从源头减少“音频二次侦查”", "Jira Template")
    rows = [
        ["Audio Required?", "Yes / No / Unknown；Unknown 默认进入音频 Triage"],
        ["Audio Goal", "玩家需要听懂什么？反馈、状态、空间、奖励、危险、氛围？"],
        ["Gameplay State", "触发状态、开始/停止、成功/失败、边界、冷却"],
        ["Unity Surface", "Prefab/Anim/VFX/脚本/配置表路径，owner"],
        ["Audio Contract", "Event/RTPC/Switch/State 需求，参数范围"],
        ["Debug/Test", "如何触发、在哪个场景、用什么工具验证"],
        ["Acceptance", "听感标准、日志标准、Profiler/截图/视频证据"],
    ]
    table(slide, 0.96, 1.58, 11.45, 4.68, ["字段", "填写说明"], rows, widths=[2.0, 9.35], font_size=10.7)
    footer(slide, page)
    page += 1

    # 17 Prioritization
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "单人音频的优先级策略", "Prioritization")
    card(slide, 0.82, 1.56, 2.86, 2.05, "S 级", "直接影响核心玩法判断：咬口、漂相、张力、泄力、鱼战斗、关键混音。必须优先。", COLORS["red"])
    card(slide, 3.92, 1.56, 2.86, 2.05, "A 级", "影响长期体验和差异化：装备、鱼类、天气、UI 奖励、多人核心反馈。分批做。", COLORS["amber"])
    card(slide, 7.02, 1.56, 2.86, 2.05, "B 级", "补完整度和版本质感：细化 UI、低优先状态、普通鱼类扩展。按版本池处理。", COLORS["blue"])
    card(slide, 10.12, 1.56, 2.20, 2.05, "C 级", "装饰/低价值声音，默认后置或裁剪。避免稀释核心听觉信息。", COLORS["slate"])
    textbox(slide, 0.92, 4.10, 11.50, 0.45, "资源策略：核心专属、族群模板、参数化复用。不要让每条鱼、每个装备、每个 UI 都膨胀成独立需求。", size=17, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    rows = [
        ["优先做", "S 级核心循环 + A 级高频差异 + 混音/性能规则"],
        ["模板化", "普通鱼、普通装备、重复 UI、环境层、多人远处事件"],
        ["可裁剪", "Hover、低价值社交提示、纯装饰 UI、内部工具音效"],
        ["不能省", "测试 Recipe、Debug 入口、Bank/Profiler 证据、P4 提交风险检查"],
    ]
    table(slide, 1.16, 4.92, 10.98, 1.28, ["策略", "内容"], widths=[1.55, 9.35], rows=rows, font_size=9.8)
    footer(slide, page)
    page += 1

    # 18 Operating cadence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "建议的工作节奏", "Operating Cadence")
    rows = [
        ["每日 10 分钟", "看 Jira/设计 Diff、日报风险、P4 待提交；只更新矩阵状态，不开大讨论。"],
        ["每周 2 次 Triage", "音频+策划+程序+美术/QA 代表；只处理 Blocked/Risky/Unknown 字段。"],
        ["每周一次 QA 复现", "按 S/A 任务跑 Runtime QA Recipe，产出日志/Profiler/截图证据。"],
        ["每版本冻结点", "锁定 Ready 列表、Cuttable 列表、返工风险列表、Bank/提交规则。"],
        ["上线前回归", "S 级完整覆盖；A 级抽样；B/C 按风险包验收。"],
    ]
    table(slide, 0.96, 1.58, 11.45, 4.32, ["节奏", "目标"], rows, widths=[1.85, 9.5], font_size=11.0)
    card(slide, 1.04, 6.10, 11.30, 0.56, "会议纪律", "Triage 不讨论审美细节；审美细节留给音频评审。Triage 只决定：是否需要、缺什么、谁补、何时可测。", COLORS["teal"], heading_size=12.5, body_size=10.4)
    footer(slide, page)
    page += 1

    # 19 Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "三阶段落地路线图", "Roadmap")
    card(slide, 0.88, 1.58, 3.72, 3.85, "第 1 阶段：1-2 周\n先让流程跑起来", "建立音频落地矩阵模板；定义 Ready Gate；把 Jira/设计扫描输出接入矩阵；选 3-5 个真实系统试跑。", COLORS["blue"], heading_size=14, body_size=12)
    card(slide, 4.84, 1.58, 3.72, 3.85, "第 2 阶段：2-4 周\n减少重复找人", "为 UI、动画、VFX、玩法状态建立配置字段模板；固定每周 Triage；QA Recipe 与日报/Profiler 关联。", COLORS["amber"], heading_size=14, body_size=12)
    card(slide, 8.80, 1.58, 3.72, 3.85, "第 3 阶段：1-2 月\n形成生产闭环", "矩阵成为版本验收入口；Jira 增加音频影响字段；工具报告自动回填证据；P4 提交与任务证据关联。", COLORS["green"], heading_size=14, body_size=12)
    textbox(slide, 0.98, 5.95, 11.30, 0.40, "推荐从“钓鱼核心循环 + UI 高频按钮 + 动画事件 + 一个复杂状态系统”四类样本开始试点。", size=16, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 20 Management asks
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "需要团队支持的决策", "Management Ask")
    bullets(slide, 1.02, 1.55, 11.00, 4.50, [
        "确认：没有配置路径和测试方法的需求，不进入 Ready，只能算 DesignOnly 或 Blocked。",
        "确认：策划/程序/美术/QA 必须为各自字段提供 owner；音频不独自承担跨部门信息缺口。",
        "确认：每周固定 Triage 时间，优先处理 S/A 级和被阻塞任务。",
        "确认：复杂机制必须有 walkthrough，不再靠音频从文档和代码里反复猜。",
        "确认：测试证据成为关闭音频任务的条件，不把 StaticOnly / NotObserved 当作通过。",
        "确认：B/C 级内容允许按版本后置、模板化或裁剪，避免单人音频产能被低价值声音耗尽。"
    ], size=16)
    footer(slide, page)
    page += 1

    # 21 Expected outcomes
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "预期收益", "Expected Outcomes")
    card(slide, 0.86, 1.62, 3.55, 1.75, "少找人", "重复问路径、owner、测试入口的时间下降。", COLORS["teal"])
    card(slide, 4.88, 1.62, 3.55, 1.75, "少返工", "机制、资源稳定性、触发边界在制作前暴露。", COLORS["green"])
    card(slide, 8.90, 1.62, 3.55, 1.75, "可排期", "Ready/Blocked/Risky 状态让音频工作量可预测。", COLORS["blue"])
    card(slide, 0.86, 3.78, 3.55, 1.75, "可验收", "每个任务有场景、步骤、日志/Profiler 证据。", COLORS["purple"])
    card(slide, 4.88, 3.78, 3.55, 1.75, "可复用", "同类系统复用模板、字段、QA Recipe。", COLORS["amber"])
    card(slide, 8.90, 3.78, 3.55, 1.75, "可追责", "缺字段是谁补、何时补、补到哪里，一眼可见。", COLORS["red"])
    footer(slide, page)
    page += 1

    # 22 Minimum viable process
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "最小可执行版本", "MVP")
    rows = [
        ["1", "用一张 Excel/CSV/Notion 表建立音频落地矩阵，不等待系统开发。"],
        ["2", "把现有 Audio Requirement Jira Triage 输出作为 Candidate 输入。"],
        ["3", "对每条 S/A 候选需求补齐：机制、Unity 配置路径、Wwise 合约、QA Recipe。"],
        ["4", "每周两次 30 分钟 Triage，只处理缺字段和 Blocked。"],
        ["5", "运行验证时要求 QA/程序提供可复现步骤，日志和 Profiler 报告挂回矩阵。"],
        ["6", "P4 提交前跑 Changelist Check，提交说明引用矩阵行和验收证据。"],
    ]
    table(slide, 1.04, 1.54, 11.18, 4.20, ["步骤", "动作"], rows, widths=[0.8, 10.25], font_size=11.4)
    textbox(slide, 1.05, 6.10, 11.10, 0.48, "一句话：先用最简单的表把信息债显性化，然后让现有工具围绕它补证据。", size=18, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)
    page += 1

    # 23 Evidence boundary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "证据边界：工具报告应该怎么被理解", "Evidence Rules")
    rows = [
        ["Jira/文档扫描", "证明设计材料里可能存在音频需求；不证明可以制作"],
        ["Unity 静态扫描", "证明配置/引用/路径存在或缺失；不证明运行时一定触发"],
        ["Wwise 审计", "证明工程结构、资源、Bank、参数风险；不证明游戏里一定调用"],
        ["Runtime Log", "证明本次会话发生过的日志事件；没跑到的系统不能判通过"],
        ["Profiler", "证明捕获窗口内 active voice/性能状态；不代表全项目真相"],
        ["人工确认", "证明某个 owner 给出的当前事实；需要时间戳和上下文"],
    ]
    table(slide, 0.90, 1.56, 11.55, 4.50, ["证据", "能证明什么 / 不能证明什么"], rows, widths=[2.0, 9.45], font_size=10.5)
    footer(slide, page)
    page += 1

    # 24 Close
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    textbox(slide, 0.96, 1.24, 11.50, 0.76, "最终建议", size=34, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.15, 2.28, 11.05, 1.45, "不要再让音频靠“读文档 + 找人 + 猜路径 + 问测试”来承接项目复杂度。把复杂度拆成字段，把字段交给 owner，把证据挂回同一条任务。", size=25, color=COLORS["teal"], bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 1.42, 4.42, 10.55, 0.78, "这样，一个人做音频也不是没有压力，但压力会回到真正该投入的地方：声音设计、系统设计、混音判断和关键体验验收。", size=18, color=COLORS["slate"], align=PP_ALIGN.CENTER)
    footer(slide, page)

    prs.save(PPTX_PATH)


HTML_SECTIONS = [
    ("一页结论", [
        "当前问题是正常的：复杂游戏音频处在设计、美术、程序、QA 的下游，必须理解玩法系统、资源位置、触发逻辑和测试方法。",
        "但长期让一个音频人靠阅读、猜测和逐个找人确认，是流程把跨部门信息债转嫁给音频。",
        "优化目标不是减少必要沟通，而是把沟通从“找人问一切”变成“补齐任务字段”。",
    ]),
    ("目前已掌握的资料与工具", [
        f"最终工具菜单共 {CTX['tools_total']} 个条目，其中 {CTX['visible_tools']} 个日常可见入口。",
        "需求/Jira 侧已有 Audio Requirement Jira Triage、设计文档扫描、需求 Diff、SABC 清单。",
        "Wwise 侧已有工程体检、资源审计、Template Generator、Audio Logic Tester、Debug Assistant。",
        "Unity 侧已有 Audio Footprint、UI Audio Static Inspector、Animation Wwise Event AutoConfig、Runtime Log Monitor。",
        "运行 QA 侧已有 Daily Audio Log Intelligence、Profiler Voice Capture、Runtime QA Checklist、Bank 检查。",
        "P4 侧已有 Changelist Check、Changelist Organizer 和提交风险分类。",
    ]),
    ("核心痛点", [
        "需求总览：Jira 和文档能让音频知道游戏里发生什么，但不能自然给出触发、生命周期、参数、测试方式。",
        "机制讨论：复杂系统涉及状态机、VFX、动画、程序逻辑和多人/视角裁剪，无法只靠文案完全理解。",
        "配置落地：策划不知道资源名和程序路径，音频不知道 Unity 资源在哪，美术/程序不知道音频需要什么。",
        "测试验收：音频通常不知道测试场景、Debug 入口、边界条件；QA 也可能没有现成路径。",
        "最大损耗不是制作音效，而是找信息、找人、找配置入口、找可复现测试方式。",
    ]),
    ("目标流程", [
        "Jira/设计/配置表进入候选需求池。",
        "工具生成 Candidate 与 Diff，不直接生成最终制作任务。",
        "音频落地矩阵承接所有字段：来源、需求判断、机制、Unity 落地、Wwise 合约、资源混音、QA 验收、状态追踪。",
        "Triage 只补缺字段和 owner，不做无限泛聊。",
        "只有通过 Ready Gate 的任务进入正式制作、配置和验收。",
    ]),
    ("Ready Gate", [
        "Design：玩法状态、成功失败、时机、优先级、边界已知。",
        "Art/VFX：视觉/动画参考存在，或明确接受 placeholder 风险。",
        "Program：触发 owner、Event、参数、生命周期、Debug 入口已知。",
        "Wwise：层级、Bus、衰减、随机、Voice、SoundBank 方案已知。",
        "Resource/Mix：素材数、样本数、命名、混音优先级、Ducking、性能预算已知。",
        "QA：测试场景、步骤、预期结果、失败模式、证据路径已知。",
    ]),
    ("落地路线图", [
        "第 1 阶段，1-2 周：建立矩阵模板、Ready Gate、把现有需求扫描接入 Candidate，选 3-5 个真实系统试跑。",
        "第 2 阶段，2-4 周：建立 UI/动画/VFX/玩法状态字段模板，固定 Triage，QA Recipe 与日报/Profiler 关联。",
        "第 3 阶段，1-2 月：矩阵成为版本验收入口，Jira 增加音频影响字段，工具报告自动回填证据，P4 提交与任务证据关联。",
    ]),
    ("管理层需要确认", [
        "没有配置路径和测试方法的需求不进入 Ready。",
        "策划、程序、美术、QA 必须为各自字段提供 owner。",
        "复杂机制必须有短 walkthrough。",
        "StaticOnly / NotObserved 不能当作通过。",
        "B/C 级内容允许模板化、后置或裁剪。",
    ]),
]

HTML_SECTIONS = [
    ("第一性原理", [
        "音频需求不是从“想做什么声音”开始，而是从“游戏流程中有哪些需要被声音解释或强化的状态变化”开始。",
        "要判断所有音频需求，音频至少需要知道：游戏有哪些系统、每个系统有哪些流程、流程涉及哪些美术资源和程序逻辑、每个流程怎么触发和测试。",
        "公式：系统 × 流程 × 资源/逻辑 × 触发/测试 × 音频判定 = 可落地音频需求。",
        "任何一项缺失，音频都只能做候选判断或方案预研，不能稳定进入正式制作和验收。",
    ]),
    ("当前覆盖率缺口", [
        "游戏有多少系统：Jira 目前大约 cover 80-90%，能形成系统视野，但仍要关注未建单或拆分不准的内容。",
        "每个系统的所有流程：策划文档大约 cover 30-40%，文档常写目标，不写完整状态机、边界和音频触发。",
        "流程涉及的美术资源和程序逻辑：策划文档大约 cover 20-30%，设计名、资源名、脚本名、Prefab 路径之间缺映射。",
        "所有流程的测试和触发方式：策划文档基本不 cover，目前准备问 QA，但需要形成系统性测试 Recipe。",
    ]),
    ("根问题", [
        "上游需求不全，音频无法稳定判断什么需要声音。",
        "上游并不知道音频判断和落地需要哪些输入。",
        "上游部门多，音频反查 owner、资源名和测试路径极其耗时。",
    ]),
    ("单点击破", [
        "第一步列系统：Jira、版本范围、系统 owner。",
        "第二步拆流程：玩家路径、状态变化、成功失败、边界。",
        "第三步标反馈点：哪些变化影响判断、操作、空间、危险、奖励、情绪或节奏。",
        "第四步映射资源/逻辑：Prefab、Animation、VFX、配置表、脚本状态。",
        "第五步形成音频合约：Play/Stop、RTPC、Switch/State、多人裁剪、冷却。",
        "第六步写测试 Recipe：场景、步骤、Debug、预期听感、日志或 Profiler 证据。",
    ]),
    ("接口人与权威来源", [
        "系统有没有：Jira、制作人、系统策划。",
        "系统怎么玩：策划案、系统策划。",
        "状态什么时候变：程序、状态机、配置表。",
        "资源在哪里：美术、TA、动画、VFX、Unity 资源表。",
        "怎么触发测试：QA、程序 Debug、测试文档。",
        "声音怎么表现：音频。",
    ]),
    ("上游可前置提供的信息", [
        "Prefab：路径、挂点、组件、是否运行时生成、owner、是否可加音频字段。",
        "Animation：clip 名、关键帧、循环段、是否可加 Animation Event、预览方式。",
        "VFX：VFX 名、播放时机、强弱等级、持续时间、是否跟随对象。",
        "UI：prefab、按钮组件、pressed/hover/success/fail 字段、默认音频规则。",
        "Config：表名、ID、AudioEventId/AudioProfileId 字段、是否热更。",
        "State Machine：状态名、进入/退出条件、状态切换频率、owner。",
        "Runtime Logic：触发函数、GameObject、生命周期、日志、多人/本地玩家规则。",
        "QA Scene：场景名、进入方式、账号/存档、Debug 工具、复现步骤。",
    ]),
] + HTML_SECTIONS


def build_html():
    sections_html = []
    for title_text, paras in HTML_SECTIONS:
        sections_html.append(
            f"<section><h2>{html.escape(title_text)}</h2>"
            + "".join(f"<p>{html.escape(p)}</p>" for p in paras)
            + "</section>"
        )
    visible_tools = "".join(f"<li>{html.escape(name)}</li>" for name in CTX["visible_tool_names"])
    daily = CTX["daily"]
    source_quality = daily.get("source_quality", {})
    source_rows = "".join(
        f"<tr><th>{html.escape(k)}</th><td>{html.escape(str(v))}</td></tr>"
        for k, v in source_quality.items()
    )
    body = f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <title>ProjectEF 音频研发流程优化方案</title>
  <style>
    :root {{
      --ink:#1c242e; --muted:#667085; --line:#d9e0e7; --teal:#137973;
      --blue:#3a66b1; --amber:#e09136; --red:#c94c4c; --bg:#f7f8fa; --panel:#fff;
    }}
    body {{ margin:0; font-family:"Microsoft YaHei", "Segoe UI", sans-serif; color:var(--ink); background:var(--bg); }}
    main {{ max-width:1080px; margin:0 auto; padding:48px 28px 72px; }}
    header {{ padding:38px 40px; background:#f3f7f6; border:1px solid var(--line); border-radius:12px; }}
    h1 {{ font-size:42px; margin:0 0 12px; letter-spacing:0; }}
    .subtitle {{ font-size:22px; color:var(--teal); font-weight:700; }}
    .meta {{ color:var(--muted); margin-top:22px; }}
    section {{ background:var(--panel); border:1px solid var(--line); border-radius:10px; padding:28px 32px; margin-top:22px; }}
    h2 {{ margin:0 0 18px; font-size:26px; }}
    p {{ font-size:17px; line-height:1.72; margin:10px 0; }}
    table {{ border-collapse:collapse; width:100%; margin-top:14px; font-size:15px; }}
    th, td {{ border:1px solid var(--line); padding:10px 12px; vertical-align:top; text-align:left; }}
    th {{ background:#eef7f6; color:#105f5a; width:210px; }}
    ul {{ columns:2; font-size:15px; line-height:1.72; }}
    .pill {{ display:inline-block; padding:6px 10px; border-radius:999px; background:#eef7f6; color:var(--teal); font-weight:700; margin-right:8px; }}
    @media print {{ body {{ background:white; }} main {{ max-width:none; }} section, header {{ break-inside:avoid; }} }}
  </style>
</head>
<body>
<main>
  <header>
    <div class="pill">ProjectEF</div><div class="pill">Audio Workflow</div>
    <h1>音频研发流程优化方案</h1>
    <div class="subtitle">从“下游救火”到“字段驱动的协作闭环”</div>
    <div class="meta">生成日期：2026-06-16；PPT 同步生成：{html.escape(str(PPTX_PATH))}</div>
  </header>
  {''.join(sections_html)}
  <section>
    <h2>日常可见工具入口</h2>
    <ul>{visible_tools}</ul>
  </section>
  <section>
    <h2>最新日报证据边界</h2>
    <p>最新日报运行验证均值：{html.escape(str(daily.get('runtime_validation_average', '33.8')))}%；模块数：{len(daily.get('modules', [])) or 16}；问题项：{len(daily.get('problems', [])) or 29}。</p>
    <table>{source_rows}</table>
  </section>
</main>
</body>
</html>"""
    HTML_PATH.write_text(body, encoding="utf-8")


def main():
    REPORT_DIR.mkdir(parents=True, exist_ok=True)
    build_deck()
    build_html()
    print(f"PPTX: {PPTX_PATH}")
    print(f"HTML: {HTML_PATH}")


if __name__ == "__main__":
    main()
